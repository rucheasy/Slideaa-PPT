from dotenv import load_dotenv
import os

load_dotenv()
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import openai
import requests
import os

app = Flask(__name__)
CORS(app)

openai.api_key = os.getenv("OPENAI_API_KEY")
# 🔹 Generate content
def generate_content(topic, slides):
    prompt = f"""
    Create {slides} slides for "{topic}".
    Format:
    Title: ...
    Points:
    - ...
    - ...
    """

    res = openai.ChatCompletion.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}]
    )

    return res.choices[0].message.content


# 🔹 Preview route
""" @app.route("/preview", methods=["POST"])
def preview():
    data = request.get_json()
    topic = data["topic"]
    slides = data["slides"]

    content = generate_content(topic, slides)

    slides_data = []
    raw = content.split("Title:")

    for s in raw[1:]:
        lines = s.strip().split("\n")
        title = lines[0]
        points = [l.replace("- ", "") for l in lines if "- " in l]

        slides_data.append({
            "title": title,
            "points": points,
            "image": ""
        })

    return jsonify({"slides": slides_data})

 """

from flask import jsonify, request

@app.route("/preview", methods=["POST"])
def preview():
    try:
        data = request.get_json()

        topic = data.get("topic", "")
        slides = int(data.get("slides", 5))

        content = generate_content(topic, slides)

        slides_data = []

        # Split slides safely
        raw_slides = content.split("Title:")

        for s in raw_slides[1:]:
            lines = s.strip().split("\n")

            if not lines:
                continue

            title = lines[0].strip()

            points = []
            for l in lines:
                if l.strip().startswith("-"):
                    points.append(l.replace("- ", "").strip())

            slides_data.append({
                "title": title,
                "points": points,
                "image": ""
            })

        # Fallback if parsing fails
        if not slides_data:
            slides_data = [{
                "title": "Sample Slide",
                "points": ["Point 1", "Point 2"],
                "image": ""
            }]

        return jsonify({"slides": slides_data})

    except Exception as e:
        print("ERROR:", e)
        return jsonify({
            "slides": [{
                "title": "Error",
                "points": ["Something went wrong"],
                "image": ""
            }]
        })
# 🔹 Template styling
def apply_template(slide, title_shape, template):

    if template == "minimal":
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255,255,255)

    elif template == "tech":
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(20,20,40)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,255,255)

    elif template == "aesthetic":
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255,228,225)
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,105,180)


# 🔹 Animation (basic)
def apply_animation(slide, animation):
    try:
        if animation == "fade":
            slide.slide_show_transition.type = 1
        elif animation == "zoom":
            slide.slide_show_transition.type = 2
        elif animation == "slide":
            slide.slide_show_transition.type = 3
    except:
        pass


# 🔹 Generate final PPT
@app.route("/generate", methods=["POST"])
def generate():
    data = request.get_json()
    slides_data = data["slides"]
    template = data.get("template", "minimal")
    animation = data.get("animation", "fade")

    prs = Presentation()

    for s in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        title = s["title"]
        points = s["points"]
        image = s.get("image")

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        title_box.text_frame.text = title

        # Content
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(4))
        tf = box.text_frame
        tf.clear()

        for p in points:
            para = tf.add_paragraph()
            para.text = p

        # Image
        if image:
            try:
                if image.startswith("http"):
                    img = requests.get(image).content
                    with open("temp.jpg", "wb") as f:
                        f.write(img)
                    img_path = "temp.jpg"
                else:
                    img_path = image

                slide.shapes.add_picture(img_path, Inches(5.5), Inches(1.5), width=Inches(4))
            except:
                pass

        apply_template(slide, title_box, template)
        apply_animation(slide, animation)

    file_path = "final.pptx"
    prs.save(file_path)
    return send_file(file_path, as_attachment=True)


if __name__ == "__main__":
    app.run(debug=True)